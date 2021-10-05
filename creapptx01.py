# pip install python-pptx
# https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61

import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR


def cabecera():
    # Creación de la ppt
    print("")   
    print("************************************************************************")
    print("*** Generador automático de plantilla de resentaciones en PowerPoint ***")
    print("************************************************************************")
    print("python-pptx v0.6.20 disponible en https://python-pptx.readthedocs.io/en/latest/\n")
    return 0

def diapositiva (left, top, logo_path, patron, titulo, subtitulo):
    # Las diapositivas se añaden en el layout
    lyt = prs.slide_layouts[patron] 
    print("[OK]\t Creado prs.slide_layouts\t\t\tslide_layouts[", patron, "]")

    # Diapositiva slide
    slide = prs.slides.add_slide(lyt)
    print("[OK]\t Añadida diapositiva\t\t\t\tprs.slides.add_slide[", patron, "]")

    # Logo
    logo = slide.shapes.add_picture(logo_path,left,top)
    print("[OK]\t Añadido slide.shapes.add_picture\t\t[",logo_path,"]", left.inches, top.inches)

    # Hay 11 patrones 0-10
    if patron > 10:
        print("[EZ]\t Layout mayor que 10")
    elif patron == 6: # slide en blanco sin placeholders
        print("[EX]\t Layout 6 sin título")
        print("[EX]\t Layout 6 sin subtítulo")
    else:
        title = slide.shapes.title
        print("[OK]\t Creado Título1 slide.shapes.title\t\t", title.name)
        title.text=titulo
        print("[OK]\t Título1 \t\t\t\t\t", title.text)
        if patron == 5:
            subtitle=slide.placeholders[0]
            print("[EX]\t Layout 5 con placeholder 0 en lugar de placeholder 1")
        else:
            subtitle=slide.placeholders[1]
            print("[OK]\t Creado Subtítulo2 slide.placeholders[1]\t", subtitle.name)
            subtitle.text=subtitulo
            print("[OK]\t Título1 \t\t\t\t\t", subtitle.text)
    return 0


def plantilla_diapositivas():
    # Añadir slides
    diapositiva (left, top, logo_path, 0, "0 Title 1", "Subtítle 2")
    diapositiva (left, top, logo_path, 1, "1 Title 1", "Content Placeholder 2")
    diapositiva (left, top, logo_path, 2, "2 Title 1", "Text Placeholder 2")
    diapositiva (left, top, logo_path, 3, "3 Title 1", "Content Placeholder 2")
    diapositiva (left, top, logo_path, 4, "4 Title 1", "Text Placeholder 2")
    diapositiva (left, top, logo_path, 5, "5 Título 5", "Placeholder 0")
    diapositiva (left, top, logo_path, 6, "6 En blanco", "En blanco")
    diapositiva (left, top, logo_path, 7, "7 Title 1", "Content Placeholder 2")
    diapositiva (left, top, logo_path, 8, "8 Title 1", "Picture Placeholder 2")
    diapositiva (left, top, logo_path, 9, "9 Title 1", "Vertical Text Placeholder 2")
    diapositiva (left, top, logo_path, 10, "10 Vertical Title 1","Vertical Text Placeholder 2")

    # Slide en blanco
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    logo=slide.shapes.add_picture(logo_path,left,top)
    return 0


# **********************
# * PROGRAMA PRINCIPAL *
# **********************

# Inicialización
left = Inches(8)
top = Inches(0.2)
logo_path = 'Logo-grupo-ITE_HR.png'
patron = 0
titulo = "Título"
subtitulo = "Subtítulo"

# Cabecera
cabecera()

# Presentación
prs = Presentation() # No olvidar poner el nombre en prs.save(prs_path)
#prs.slide_width = Inches(16)
#prs.slide_height = Inches(9)



# IV Curso de Inteligencia y Contrainteligencia

# Slide 3 INTRODUCCIÓN: Título
titulo="IV Curso de Inteligencia y Contrainteligencia"
subtitulo="Módulo V. Ciberinteligencia: Una mirada al Ciberespionaje"
diapositiva (left, top, logo_path, 0, titulo, subtitulo)

# Slide 2 INTRODUCCIÓN: Objetivos 
lyt = prs.slide_layouts[1] 
slide = prs.slides.add_slide(lyt)
shapes = slide.shapes
logo = slide.shapes.add_picture(logo_path,left,top)
title_shape = slide.shapes.title
body_shape = slide.shapes.placeholders[1]
titulo="INTRODUCCIÓN"
title_shape.text = titulo

tf = body_shape.text_frame
tf.text = "OBJETIVOS"
p = tf.add_paragraph()
p.text = "Los objetivos del curso son los siguientes"
p.level = 1
p = tf.add_paragraph()
p.text = "Conoce las técnicas para obtener y analizar información para identificar y predecir amenazas y ataques cibernéticos, sobre todo para diseñar acciones y políticas de prevención"
p.level = 2
p = tf.add_paragraph()
p.text = "La ciberdelincuencia evolucionó aún más durante la pandemia, hay nuevas amenazas y formas en las que nuestra información se ve amenazada en la reds"
p.level = 2
p = tf.add_paragraph()
p.text = "También analizaremos como desarrollar capacidades de Ciber inteligencia que fortalezcan las estrategias de protección, prevención y respuesta."
p.level = 2


# Slide 3 INTRODUCCIÓN: Audiencia 
lyt = prs.slide_layouts[1] 
slide = prs.slides.add_slide(lyt)
shapes = slide.shapes
logo = slide.shapes.add_picture(logo_path,left,top)
title_shape = slide.shapes.title
body_shape = slide.shapes.placeholders[1]
titulo="INTRODUCCIÓN"
title_shape.text = titulo

tf = body_shape.text_frame
tf.text = "AUDIENCIA"
p = tf.add_paragraph()
p.text = "Este curso está dirigodo a los siguientes perfiles"
p.level = 1
p = tf.add_paragraph()
p.text = "Funcionarios de Inteligencia"
p.level = 2
p = tf.add_paragraph()
p.text = "Asesores Políticos"
p.level = 2
p = tf.add_paragraph()
p.text = "Consultores de Seguridad"
p.level = 2
p = tf.add_paragraph()
p.text = "Analistas de Inteligencia"
p.level = 2
p = tf.add_paragraph()
p.text = "Investigadores privados"
p.level = 2
p = tf.add_paragraph()
p.text = "Científicos afines a la Inteligencia y ciberinteligencia"
p.level = 2



# Slide 4 HISTORIA 1



# Slide 5 HISTORIA 2



# Slide 6 HISTORIA 3



# Slide 7 HISTORIA 4



# Slide 8 HISTORIA 5 



# Plantila de diapositivas
plantilla_diapositivas()

# Grabación de la ppt
prs_path = "ITE_Plantilla_Presentación_automática.pptx"
prs.save(prs_path) # saving file
print("[OK]\t Presentación generada\t\t\t\t", prs_path,)